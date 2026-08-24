"""Phase 1 parsing tests with inline fixtures."""

from __future__ import annotations

import io
import json

import pymupdf as fitz
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.graph.state import BlockType, SourceFormat
from app.parsing.html_parser import parse_html
from app.parsing.ipynb_parser import build_sample_ipynb, build_sample_ipynb_with_output, parse_ipynb
from app.parsing.normalize import combine_text, normalize_documents
from app.parsing.pdf_parser import parse_pdf
from app.parsing.pptx_parser import parse_pptx
from app.parsing.text_parser import parse_text


def _make_pdf(text: str = "Hello PDF world.\n\nSecond paragraph here.") -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text, fontsize=12)
    return doc.tobytes()


def _make_pptx(slides: list[str] | None = None) -> bytes:
    slides = slides or ["Slide One Title", "Slide Two Content"]
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in slides:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tx.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_html(body: str = "<html><body><h1>Title</h1><p>Paragraph one.</p><p>Paragraph two.</p></body></html>") -> bytes:
    return body.encode("utf-8")


SAMPLE_TRANSCRIPT = b"""Speaker A: Welcome to the session on AI agents.

Speaker B: Today we cover LangGraph and multi-agent systems.

Speaker A: Let's start with the basics of stateful graphs.
"""


class TestPDFParser:
    def test_pdf_block_count_and_chars(self):
        raw = "Alpha beta gamma.\n\nDelta epsilon zeta."
        content = _make_pdf(raw)
        doc = parse_pdf(content, "sample.pdf", "pdf1")
        assert doc.source_format == SourceFormat.PDF
        assert doc.block_count >= 1
        assert doc.char_count >= len(raw.replace("\n\n", " ")) * 0.5  # tolerance for extraction

    def test_pdf_page_metadata(self):
        doc = parse_pdf(_make_pdf("Page content"), "sample.pdf", "pdf1")
        assert all(b.page_or_slide is not None for b in doc.blocks)


class TestPPTXParser:
    def test_pptx_slides(self):
        content = _make_pptx(["First slide", "Second slide"])
        doc = parse_pptx(content, "deck.pptx", "pptx1")
        assert doc.source_format == SourceFormat.PPTX
        assert doc.block_count >= 2
        combined = " ".join(b.text for b in doc.blocks)
        assert "First slide" in combined
        assert "Second slide" in combined
        assert any(b.block_type == BlockType.SLIDE for b in doc.blocks)

    def test_pptx_tables_and_notes(self):
        prs = Presentation()
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(1.2))
        table_shape.table.cell(0, 0).text = "Metric"
        table_shape.table.cell(0, 1).text = "Value"
        table_shape.table.cell(1, 0).text = "Vocab"
        table_shape.table.cell(1, 1).text = "50257"
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Remember the GPT-2 vocabulary size."
        buf = io.BytesIO()
        prs.save(buf)
        doc = parse_pptx(buf.getvalue(), "table.pptx", "pptx2")
        types = {b.block_type for b in doc.blocks}
        assert BlockType.TABLE in types
        assert "50257" in " ".join(b.text for b in doc.blocks)
        assert any("Speaker notes" in b.text for b in doc.blocks)


class TestHTMLParser:
    def test_html_extraction(self):
        doc = parse_html(_make_html(), "page.html", "html1")
        assert doc.source_format == SourceFormat.HTML
        assert doc.block_count >= 1
        combined = " ".join(b.text for b in doc.blocks)
        assert "Paragraph" in combined or "Title" in combined

    def test_html_keeps_headings_lists_tables_code(self):
        html = b"""<html><body>
        <nav>Skip me</nav>
        <h1>BPE</h1>
        <p>Byte Pair Encoding merges frequent pairs.</p>
        <ul><li>Start from characters</li><li>Merge the winner</li></ul>
        <pre>print(50257)</pre>
        <table><tr><th>Model</th><th>Vocab</th></tr><tr><td>GPT-2</td><td>50257</td></tr></table>
        </body></html>"""
        doc = parse_html(html, "bpe.html", "html2")
        types = {b.block_type for b in doc.blocks}
        assert BlockType.HEADING in types
        assert BlockType.LIST in types
        assert BlockType.CODE in types
        assert BlockType.TABLE in types
        assert "50257" in " ".join(b.text for b in doc.blocks)


class TestIPYNBParser:
    def test_ipynb_cells(self):
        content = build_sample_ipynb()
        doc = parse_ipynb(content, "nb.ipynb", "nb1")
        assert doc.source_format == SourceFormat.IPYNB
        assert doc.block_count == 2
        types = {b.block_type for b in doc.blocks}
        assert BlockType.HEADING in types or BlockType.PARAGRAPH in types
        assert BlockType.CODE in types

    def test_ipynb_keeps_cell_outputs(self):
        doc = parse_ipynb(build_sample_ipynb_with_output(), "out.ipynb", "nb2")
        texts = " ".join(b.text for b in doc.blocks)
        assert "print(50257)" in texts
        assert "50257" in texts
        assert any(b.metadata.get("kind") == "output" for b in doc.blocks)


class TestTranscriptParser:
    def test_transcript_paragraphs(self):
        doc = parse_text(SAMPLE_TRANSCRIPT, "session.transcript", "tx1")
        assert doc.source_format == SourceFormat.TRANSCRIPT
        assert doc.block_count >= 2
        assert doc.char_count > 50


class TestNormalize:
    def test_global_order_unique(self):
        docs = [
            parse_pdf(_make_pdf("PDF text"), "a.pdf", "pdf1"),
            parse_text(SAMPLE_TRANSCRIPT, "b.transcript", "tx1"),
        ]
        normalized = normalize_documents(docs)
        orders = [b.order for doc in normalized for b in doc.blocks]
        assert orders == sorted(orders)
        assert len(orders) == len(set(orders))

    def test_block_id_format(self):
        doc = parse_pdf(_make_pdf("Test"), "a.pdf", "pdf1")
        normalized = normalize_documents([doc])
        bid = normalized[0].blocks[0].block_id
        assert bid.startswith("pdf1-")
        assert "-b" in bid

    def test_combine_text(self):
        docs = normalize_documents([
            parse_text(b"Hello world", "t.txt", "t1"),
        ])
        combined = combine_text(docs)
        assert "Hello world" in combined
        assert "t.txt" in combined
        assert "[paragraph" in combined

    def test_markdown_headings_keep_level(self):
        doc = parse_text(b"# Byte Pair Encoding\n\nA merge process.", "notes.md", "t2")
        assert doc.blocks[0].block_type == BlockType.HEADING
        assert doc.blocks[0].text == "Byte Pair Encoding"
        assert doc.blocks[0].metadata.get("heading_level") == 1


class TestSourcePack:
    def test_pack_keeps_headings_from_every_file(self):
        from app.graph.state import Block, DocumentIR, SourceFormat
        from app.parsing.source_pack import pack_source

        def doc(source_id: str, filename: str, heading: str, filler: str) -> DocumentIR:
            return DocumentIR(
                source_id=source_id,
                source_format=SourceFormat.TEXT,
                filename=filename,
                blocks=[
                    Block(
                        block_id=f"{source_id}-p0-b0",
                        source_id=source_id,
                        source_format=SourceFormat.TEXT,
                        block_type=BlockType.HEADING,
                        text=heading,
                        order=1,
                    ),
                    Block(
                        block_id=f"{source_id}-p0-b1",
                        source_id=source_id,
                        source_format=SourceFormat.TEXT,
                        block_type=BlockType.PARAGRAPH,
                        text=filler,
                        order=2,
                    ),
                ],
                char_count=len(heading) + len(filler),
                block_count=2,
            )

        filler = "x" * 8000
        docs = [
            doc("a", "one.md", "Heading from file one", filler),
            doc("b", "two.md", "Heading from file two", filler),
            doc("c", "three.md", "Heading from file three", filler),
        ]
        packed = pack_source(docs, budget=6000)
        assert "Heading from file one" in packed
        assert "Heading from file two" in packed
        assert "Heading from file three" in packed
        assert "PACKED SOURCE" in packed
