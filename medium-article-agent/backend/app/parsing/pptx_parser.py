"""PPTX parser: shapes, grouped shapes, tables, speaker notes."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.graph.state import Block, BlockType, DocumentIR, SourceFormat
from app.parsing.normalize import make_block_id


def _iter_shapes(shapes):
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def _table_markdown(shape) -> str:
    table = shape.table
    lines: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _is_title(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        ptype = shape.placeholder_format.type
        return ptype in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.VERTICAL_TITLE)
    except Exception:
        return False


def parse_pptx(content: bytes, filename: str, source_id: str) -> DocumentIR:
    prs = Presentation(io.BytesIO(content))
    blocks: list[Block] = []
    warnings: list[str] = []
    order = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        local_idx = 0
        slide_had_content = False
        slide_bits: list[str] = []

        for shape in _iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                md = _table_markdown(shape)
                if not md.strip():
                    continue
                slide_had_content = True
                order += 1
                local_idx += 1
                blocks.append(
                    Block(
                        block_id=make_block_id(source_id, slide_num, local_idx),
                        source_id=source_id,
                        source_format=SourceFormat.PPTX,
                        block_type=BlockType.TABLE,
                        text=md,
                        order=order,
                        page_or_slide=slide_num,
                        metadata={"kind": "table"},
                    )
                )
                continue

            text = ""
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
            elif hasattr(shape, "text"):
                text = (shape.text or "").strip()
            if not text:
                continue
            slide_had_content = True
            slide_bits.append(text)
            btype = BlockType.HEADING if _is_title(shape) or local_idx == 0 else BlockType.PARAGRAPH
            if "\n" in text and any(line.strip()[:1] in {"•", "-", "*"} for line in text.splitlines()):
                btype = BlockType.LIST
            order += 1
            local_idx += 1
            blocks.append(
                Block(
                    block_id=make_block_id(source_id, slide_num, local_idx),
                    source_id=source_id,
                    source_format=SourceFormat.PPTX,
                    block_type=btype,
                    text=text,
                    order=order,
                    page_or_slide=slide_num,
                    metadata={"shape": getattr(shape, "name", "")},
                )
            )

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        if notes:
            slide_had_content = True
            order += 1
            local_idx += 1
            blocks.append(
                Block(
                    block_id=make_block_id(source_id, slide_num, local_idx),
                    source_id=source_id,
                    source_format=SourceFormat.PPTX,
                    block_type=BlockType.METADATA,
                    text=f"[Speaker notes] {notes}",
                    order=order,
                    page_or_slide=slide_num,
                    metadata={"kind": "notes"},
                )
            )

        if not slide_had_content:
            warnings.append(f"slide {slide_num}: no extractable text, tables, or notes")
            order += 1
            local_idx += 1
            blocks.append(
                Block(
                    block_id=make_block_id(source_id, slide_num, local_idx),
                    source_id=source_id,
                    source_format=SourceFormat.PPTX,
                    block_type=BlockType.SLIDE,
                    text=f"[Empty slide {slide_num}]",
                    order=order,
                    page_or_slide=slide_num,
                    metadata={"empty_slide": True},
                )
            )
        elif slide_bits:
            # Keep a slide-level rollup so older tests and planners still see one unit per slide.
            order += 1
            local_idx += 1
            blocks.append(
                Block(
                    block_id=make_block_id(source_id, slide_num, local_idx),
                    source_id=source_id,
                    source_format=SourceFormat.PPTX,
                    block_type=BlockType.SLIDE,
                    text="\n".join(slide_bits),
                    order=order,
                    page_or_slide=slide_num,
                    metadata={"kind": "slide_rollup"},
                )
            )

    if not blocks:
        warnings.append("no extractable slides")
        blocks.append(
            Block(
                block_id=make_block_id(source_id, 1, 0),
                source_id=source_id,
                source_format=SourceFormat.PPTX,
                block_type=BlockType.METADATA,
                text=f"[No extractable content from {filename}]",
                order=1,
                page_or_slide=1,
            )
        )

    char_count = sum(len(b.text) for b in blocks)
    return DocumentIR(
        source_id=source_id,
        source_format=SourceFormat.PPTX,
        filename=filename,
        blocks=blocks,
        char_count=char_count,
        block_count=len(blocks),
        warnings=warnings,
    )
